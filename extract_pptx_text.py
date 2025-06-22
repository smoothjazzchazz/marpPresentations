import sys
from pptx import Presentation

# Path to the PowerPoint file
pptx_path = "Preliminary draft.pptx"

# Load the presentation
prs = Presentation(pptx_path)

for i, slide in enumerate(prs.slides, 1):
    print(f"\n--- Slide {i} ---\n")
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text:
                print(text) 